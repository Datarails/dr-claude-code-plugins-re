"""PowerPoint presentation generation utilities.

Provides templates and helpers for creating professional PowerPoint presentations
with consistent formatting, layouts, and styling.
"""

from typing import List, Dict, Tuple, Optional, Any
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime


class PowerPointReport:
    """Builder for creating PowerPoint presentations."""

    def __init__(self, title: str = "Report", subtitle: str = ""):
        """Initialize PowerPoint presentation.

        Args:
            title: Presentation title
            subtitle: Optional subtitle
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # Color scheme
        self.colors = {
            "primary": RGBColor(46, 117, 182),        # Blue
            "secondary": RGBColor(237, 125, 49),      # Orange
            "success": RGBColor(112, 173, 71),        # Green
            "warning": RGBColor(255, 192, 0),         # Yellow
            "danger": RGBColor(197, 80, 77),          # Red
            "text": RGBColor(51, 51, 51),             # Dark gray
            "light_text": RGBColor(127, 127, 127),    # Light gray
            "white": RGBColor(255, 255, 255),         # White
        }

        # Add title slide
        if title:
            self._add_title_slide(title, subtitle)

    def _add_title_slide(self, title: str, subtitle: str = ""):
        """Add title slide.

        Args:
            title: Slide title
            subtitle: Slide subtitle
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["primary"]

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.colors["white"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(9), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = self.colors["white"]
            p.alignment = PP_ALIGN.CENTER

        # Footer with date
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7), Inches(9), Inches(0.4)
        )
        footer_frame = footer_box.text_frame
        p = footer_frame.paragraphs[0]
        p.text = datetime.now().strftime("%B %d, %Y")
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = self.colors["white"]
        p.alignment = PP_ALIGN.CENTER

    def add_slide(self, title: str = "", layout_type: str = "title_and_content") -> Any:
        """Add a slide to the presentation.

        Args:
            title: Slide title
            layout_type: Layout type (title_and_content, title_only, blank)

        Returns:
            Slide object
        """
        layout_map = {
            "title_and_content": 1,
            "title_only": 5,
            "blank": 6,
        }

        layout_idx = layout_map.get(layout_type, 1)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["white"]

        # Add title if provided
        if title and layout_type != "blank":
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].font.size = Pt(40)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = self.colors["primary"]

        return slide

    def add_title_and_bullets(self, slide: Any, title: str, bullets: List[str]):
        """Add slide with title and bullet points.

        Args:
            slide: Slide object
            title: Slide title
            bullets: List of bullet points
        """
        if not slide.shapes.title:
            return

        slide.shapes.title.text = title

        # Add bullet points
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for bullet_text in bullets:
                p = text_frame.add_paragraph()
                p.text = bullet_text
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = self.colors["text"]

    def add_text_box(self, slide: Any, left: float, top: float, width: float,
                     height: float, text: str, font_size: int = 14,
                     bold: bool = False, color: Optional[RGBColor] = None):
        """Add text box to slide.

        Args:
            slide: Slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Text content
            font_size: Font size in points
            bold: Whether text is bold
            color: Text color (default: primary color)
        """
        if color is None:
            color = self.colors["text"]

        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color

    def add_image(self, slide: Any, image_path: str, left: float, top: float,
                  width: float = 5, height: float = 3.75):
        """Add image to slide.

        Args:
            slide: Slide object
            image_path: Path to image file
            left: Left position in inches
            top: Top position in inches
            width: Image width in inches
            height: Image height in inches
        """
        slide.shapes.add_picture(
            image_path,
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height)
        )

    def add_metrics_boxes(self, slide: Any, metrics: Dict[str, str], columns: int = 3):
        """Add metric boxes (KPI style) to slide.

        Args:
            slide: Slide object
            metrics: Dictionary of {metric_name: metric_value}
            columns: Number of columns
        """
        box_width = 2.5
        box_height = 1.2
        spacing_x = 3
        spacing_y = 1.5

        start_left = 0.5
        start_top = 1.5

        for idx, (name, value) in enumerate(metrics.items()):
            row = idx // columns
            col = idx % columns

            left = start_left + col * spacing_x
            top = start_top + row * spacing_y

            # Background box
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(left),
                Inches(top),
                Inches(box_width),
                Inches(box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.line.color.rgb = self.colors["primary"]
            shape.line.width = Pt(2)

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(top + 0.1),
                Inches(box_width - 0.2), Inches(0.4)
            )
            label_frame = label_box.text_frame
            p = label_frame.paragraphs[0]
            p.text = name
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.colors["light_text"]

            # Value
            value_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(top + 0.5),
                Inches(box_width - 0.2), Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.word_wrap = True
            p = value_frame.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors["primary"]

    def add_two_column_layout(self, slide: Any, left_content: Dict[str, Any],
                               right_content: Dict[str, Any]):
        """Add two-column layout to slide.

        Args:
            slide: Slide object
            left_content: Dictionary with 'title' and 'items' for left column
            right_content: Dictionary with 'title' and 'items' for right column
        """
        # Left column
        if "title" in left_content:
            self.add_text_box(slide, 0.5, 1.5, 4.5, 0.5,
                            left_content["title"], font_size=18, bold=True)

        if "items" in left_content:
            items_text = "\n".join(left_content["items"])
            self.add_text_box(slide, 0.5, 2.1, 4.5, 4,
                            items_text, font_size=14)

        # Right column
        if "title" in right_content:
            self.add_text_box(slide, 5.5, 1.5, 4, 0.5,
                            right_content["title"], font_size=18, bold=True)

        if "items" in right_content:
            items_text = "\n".join(right_content["items"])
            self.add_text_box(slide, 5.5, 2.1, 4, 4,
                            items_text, font_size=14)

    def save(self, filepath: str):
        """Save presentation to file.

        Args:
            filepath: Output file path

        Returns:
            File path
        """
        self.prs.save(filepath)
        return filepath

    def to_bytes(self) -> bytes:
        """Convert presentation to bytes.

        Returns:
            Bytes representation of presentation
        """
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
